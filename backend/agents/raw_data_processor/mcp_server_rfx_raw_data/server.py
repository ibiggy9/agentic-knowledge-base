"""Document processing MCP server for Google Drive integration"""

import os
import logging
import asyncio
from mcp.server.fastmcp import FastMCP
from google.oauth2.service_account import Credentials as ServiceAccountCredentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
import io
import tempfile
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import PyPDF2
from dotenv import load_dotenv

load_dotenv()

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger("document-processor")

mcp = FastMCP("document-processor")

document_cache = {}

SERVER_CONFIG = {
    "credentials_path": os.environ.get('CREDENTIALS_PATH'),
    "base_folder_id": os.environ.get('BASE_FOLDER_ID'),
    "base_folder_name": os.environ.get('BASE_FOLDER_NAME')
}

def get_drive_service():
    try:
        service_account_file = os.environ.get('GOOGLE_SERVICE_ACCOUNT_FILE')

        if not service_account_file or not os.path.exists(service_account_file):
            service_account_file = SERVER_CONFIG["credentials_path"]

        if os.path.exists(service_account_file):
            creds = ServiceAccountCredentials.from_service_account_file(
                service_account_file,
                scopes=['https://www.googleapis.com/auth/drive.readonly']
            )
            return build('drive', 'v3', credentials=creds)

        else:
            logger.error("GOOGLE_SERVICE_ACCOUNT_FILE environment variable not set")

        application_creds = os.environ.get('GOOGLE_APPLICATION_CREDENTIALS')
        if application_creds and os.path.exists(application_creds):
            return build('drive', 'v3', credentials=None)

        raise ValueError(
            "No Google Drive credentials found. Please set one of these environment variables:\n"
            "- GOOGLE_SERVICE_ACCOUNT_FILE: Path to service account key file (recommended)\n"
            "- GOOGLE_APPLICATION_CREDENTIALS: Path for Google Application Default Credentials\n"
            "- GOOGLE_CREDENTIALS_FILE: Path to user credentials JSON file\n"
            "- GOOGLE_CREDENTIALS: JSON string containing credentials"
        )

    except Exception as e:
        logger.error(f"Error authenticating with Google Drive: {str(e)}")
        raise

def extract_text_from_file(file_content, mime_type, file_name):
    """Enhanced file content extraction with better error handling"""
    try:
        if 'pdf' in mime_type:
            return extract_text_from_pdf(file_content)
        elif 'spreadsheet' in mime_type or file_name.endswith('.xlsx') or file_name.endswith('.xls'):
            return extract_text_from_excel(file_content)
        elif 'presentation' in mime_type or file_name.endswith(('.pptx', '.ppt')):
            return extract_text_from_presentation(file_content)
        elif 'document' in mime_type or file_name.endswith(('.docx', '.doc')):
            return extract_text_from_docx(file_content)
        elif 'text/plain' in mime_type or file_name.endswith('.txt'):
            return file_content.decode('utf-8', errors='replace')
        else:
            return f"[Unsupported file type: {mime_type} ({file_name})]"
    except Exception as e:
        logger.error(f"Error extracting text from {file_name} ({mime_type}): {str(e)}")
        return f"[Error extracting text: {str(e)}]"

def extract_text_from_pdf(file_content):
    """Enhanced PDF text extraction with better error handling"""
    text = []
    with io.BytesIO(file_content) as pdf_file:
        try:
            reader = PyPDF2.PdfReader(pdf_file)
            text.append(f"PDF Document with {len(reader.pages)} pages\n")

            for page_num in range(len(reader.pages)):
                try:
                    page_text = reader.pages[page_num].extract_text()
                    if page_text and page_text.strip():
                        text.append(f"\n--- Page {page_num + 1} ---\n")
                        text.append(page_text)
                    else:
                        text.append(f"\n--- Page {page_num + 1} (No extractable text) ---\n")
                except Exception as e:
                    text.append(f"\n--- Page {page_num + 1} (Error extracting text: {str(e)}) ---\n")
        except Exception as e:
            text.append(f"Error extracting PDF text: {str(e)}")

    return "\n".join(text)

def extract_text_from_excel(file_content):
    text_output = []

    with io.BytesIO(file_content) as file_obj:
        try:
            excel_file = pd.ExcelFile(file_obj)
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                text_output.append(f"\n\n=== SHEET: {sheet_name} ===\n")
                if df.empty:
                    text_output.append("(Empty sheet)")
                    continue

                text_output.append("COLUMNS: " + ", ".join(df.columns.astype(str)) + "\n")
                text_output.append(df.to_string(index=True))

                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) > 0:
                    text_output.append("\nNUMERIC COLUMN STATISTICS:")
                    stats = df[numeric_cols].describe().to_string()
                    text_output.append(stats)
        except Exception as e:
            text_output.append(f"Error processing Excel file: {str(e)}")

    return "\n".join(text_output)

def extract_text_from_docx(file_content):
    text = []
    with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp:
        temp.write(file_content)
        temp_path = temp.name

    try:
        doc = DocxDocument(temp_path)
        for para in doc.paragraphs:
            text.append(para.text)

        for table in doc.tables:
            text.append("\nTABLE CONTENT:")
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                text.append(" | ".join(row_text))
    except Exception as e:
        text.append(f"Error extracting DOCX text: {str(e)}")
    finally:
        try:
            os.unlink(temp_path)
        except:
            pass

    return "\n".join(text)

def extract_text_from_presentation(file_content):
    text_output = []

    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp:
        temp.write(file_content)
        temp_path = temp.name

    try:
        pres = Presentation(temp_path)

        for i, slide in enumerate(pres.slides):
            text_output.append(f"\n=== SLIDE {i+1} ===\n")

            if slide.shapes.title and slide.shapes.title.text:
                text_output.append(f"TITLE: {slide.shapes.title.text}\n")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_output.append(shape.text)

                if hasattr(shape, "has_table") and shape.has_table:
                    text_output.append("\nTABLE CONTENT:")
                    table = shape.table
                    for r in range(len(table.rows)):
                        row_content = []
                        for c in range(len(table.columns)):
                            cell = table.cell(r, c)
                            row_content.append(cell.text)
                        text_output.append(" | ".join(row_content))
    except Exception as e:
        text_output.append(f"Error processing presentation: {str(e)}")
    finally:
        try:
            os.unlink(temp_path)
        except:
            pass

    return "\n".join(text_output)

def create_document_chunks(content, chunk_size=4000, overlap=200):
    if len(content) <= chunk_size:
        return [content]
    chunks = []
    start = 0

    while start < len(content):
        end = min(start + chunk_size, len(content))

        if end < len(content):
            natural_break = content.rfind('\n', end - 500, end)
            if natural_break != -1:
                end = natural_break + 1

        chunks.append(content[start:end])
        start = end - overlap

    return chunks

async def summarize_chunk(ctx, chunk, focus_area=None, chunk_label=""):
    system_prompt = f"""You are analyzing a portion of a document ({chunk_label}).
    Provide a concise summary of the key information in this section.
    Focus on extracting the most important facts, figures, and insights."""

    if focus_area:
        system_prompt += f" Pay particular attention to content related to {focus_area}."

    user_message = f"Document section content:\n\n{chunk}"

    try:
        sampling_result = await ctx.request_context.session.create_message({
            "messages": [
                {
                    "role": "user",
                    "content": {
                        "type": "text",
                        "text": user_message
                    }
                }
            ],
            "systemPrompt": system_prompt,
            "maxTokens": 500,
            "temperature": 0.1,
            "modelPreferences": {
                "hints": [{"name": "claude-3-5-sonnet"}]
            }
        })

        if sampling_result and sampling_result.content.type == "text":
            return f"{chunk_label}:\n{sampling_result.content.text}"
        else:
            return f"{chunk_label}: Error generating summary"
    except Exception as e:
        logger.error(f"Error summarizing chunk: {str(e)}")
        return f"{chunk_label}: Error summarizing content - {str(e)}"

@mcp.tool()
async def list_drive_folders(parent_folder_id: str = None):
    """
    Lists all folders within a specified Google Drive folder.
    
    This tool retrieves all folders that are direct children of the specified parent folder.
    If no parent folder ID is provided, it uses the configured base folder.
    
    Parameters:
        parent_folder_id (str, optional): The ID of the parent folder to list folders from.
                                         Must be a valid Google Drive folder ID.
                                         If not provided, the base folder configured in SERVER_CONFIG will be used.
    
    Returns:
        list: A list of folder objects, each containing:
              - id: The unique Google Drive ID of the folder
              - name: The display name of the folder
              - mimeType: Always 'application/vnd.google-apps.folder'
              - createdTime: Timestamp when the folder was created
              - modifiedTime: Timestamp when the folder was last modified
        dict: Error response if the operation fails, with structure:
              {"error": True, "message": "Error message"}
    """
    try:
        drive_service = get_drive_service()
        actual_parent = parent_folder_id if parent_folder_id else SERVER_CONFIG["base_folder_id"]
        query = f"'{actual_parent}' in parents and mimeType='application/vnd.google-apps.folder' and trashed=false"

        response = drive_service.files().list(
            q=query,
            spaces='drive',
            fields='files(id, name, mimeType, createdTime, modifiedTime)'
        ).execute()

        folders = response.get('files', [])
        return folders
    except Exception as e:
        logger.error(f"Error listing folders: {str(e)}")
        return {"error": True, "message": f"Error listing folders: {str(e)}"}

@mcp.tool()
async def list_drive_files(folder_id: str = None, file_types: str = "all"):
    """
    Lists all files within a specified Google Drive folder with optional filtering by file type.
    
    This tool retrieves all files (non-folder items) that are direct children of the specified folder.
    If no folder ID is provided, it uses the configured base folder.
    
    Parameters:
        folder_id (str, optional): The ID of the folder to list files from.
                                  Must be a valid Google Drive folder ID.
                                  If not provided, the base folder configured in SERVER_CONFIG will be used.
        file_types (str, optional): Comma-separated list of file types to filter by.
                                   Default is "all" which returns all file types.
                                   Supported values: "pdf", "excel"/"xlsx"/"xls", "word"/"doc"/"docx",
                                   "ppt"/"pptx"/"presentation", "txt"/"text"
                                   Example: "pdf,excel,word" will return only PDF and Office documents
    
    Returns:
        list: A list of file objects, each containing:
              - id: The unique Google Drive ID of the file
              - name: The display name of the file
              - mimeType: The MIME type of the file
              - createdTime: Timestamp when the file was created
              - modifiedTime: Timestamp when the file was last modified
              - size: The size of the file in bytes (if available)
        dict: Error response if the operation fails, with structure:
              {"error": True, "message": "Error message"}
    """
    try:
        drive_service = get_drive_service()
        actual_folder = folder_id if folder_id else SERVER_CONFIG["base_folder_id"]
        query = f"'{actual_folder}' in parents and mimeType!='application/vnd.google-apps.folder' and trashed=false"

        if file_types.lower() != "all":
            type_queries = []
            types = [t.strip().lower() for t in file_types.split(",")]

            for file_type in types:
                if file_type == "pdf":
                    type_queries.append("mimeType='application/pdf'")
                elif file_type in ["excel", "xlsx", "xls"]:
                    type_queries.append("mimeType='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'")
                    type_queries.append("mimeType='application/vnd.ms-excel'")
                    type_queries.append("mimeType='application/vnd.google-apps.spreadsheet'")
                elif file_type in ["word", "doc", "docx"]:
                    type_queries.append("mimeType='application/vnd.openxmlformats-officedocument.wordprocessingml.document'")
                    type_queries.append("mimeType='application/msword'")
                    type_queries.append("mimeType='application/vnd.google-apps.document'")
                elif file_type in ["ppt", "pptx", "presentation"]:
                    type_queries.append("mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'")
                    type_queries.append("mimeType='application/vnd.ms-powerpoint'")
                    type_queries.append("mimeType='application/vnd.google-apps.presentation'")
                elif file_type in ["txt", "text"]:
                    type_queries.append("mimeType='text/plain'")

            if type_queries:
                query += " and (" + " or ".join(type_queries) + ")"

        response = drive_service.files().list(
            q=query,
            spaces='drive',
            fields='files(id, name, mimeType, createdTime, modifiedTime, size)'
        ).execute()

        files = response.get('files', [])
        return files
    except Exception as e:
        logger.error(f"Error listing files: {str(e)}")
        return {"error": True, "message": f"Error listing files: {str(e)}"}

@mcp.tool()
async def get_file_content(file_id: str, max_chars: int = 100000, include_metadata: bool = False):
    """
    Retrieves and extracts the text content from a file in Google Drive.
    
    This tool downloads the specified file from Google Drive and extracts its textual content.
    It supports various file types including PDF, Excel, Word, PowerPoint, and plain text.
    The extracted content is cached to improve performance for subsequent requests.
    
    Parameters:
        file_id (str): The ID of the file to retrieve.
                      Must be a valid Google Drive file ID.
        max_chars (int, optional): Maximum number of characters to return from the file content.
                                  Default is 100000. Set higher for larger files, but be aware of
                                  response size limitations.
        include_metadata (bool, optional): Whether to include detailed file metadata in the response.
                                          Default is False.
    
    Returns:
        dict: A dictionary containing the file content and metadata:
              - name: The display name of the file
              - mime_type: The MIME type of the file
              - content: The extracted text content (up to max_chars)
              - truncated: Boolean indicating if the content was truncated
              - total_length: Total length of the extracted content in characters
              - from_cache: Boolean indicating if the content was retrieved from cache
              - metadata: (If include_metadata=True) Additional file metadata
    """
    try:
        if file_id in document_cache:
            content, metadata = document_cache[file_id]
            result = {
                "name": metadata.get("name", "Unknown"),
                "mime_type": metadata.get("mimeType", "Unknown"),
                "content": content[:max_chars],
                "truncated": len(content) > max_chars,
                "total_length": len(content),
                "from_cache": True
            }
            
            if include_metadata:
                result["metadata"] = metadata
                
            return result

        drive_service = get_drive_service()
        file_metadata = drive_service.files().get(fileId=file_id, fields='name,mimeType,size,createdTime,modifiedTime,parents').execute()

        if file_metadata['mimeType'].startswith('application/vnd.google-apps.'):
            if file_metadata['mimeType'] == 'application/vnd.google-apps.document':
                request = drive_service.files().export_media(fileId=file_id, mimeType='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                file_content = io.BytesIO()
                downloader = MediaIoBaseDownload(file_content, request)

                done = False
                while not done:
                    status, done = downloader.next_chunk()

                content_bytes = file_content.getvalue()
                extracted_text = extract_text_from_docx(content_bytes)
            elif file_metadata['mimeType'] == 'application/vnd.google-apps.spreadsheet':
                request = drive_service.files().export_media(fileId=file_id, mimeType='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
                file_content = io.BytesIO()
                downloader = MediaIoBaseDownload(file_content, request)

                done = False
                while not done:
                    status, done = downloader.next_chunk()

                content_bytes = file_content.getvalue()
                extracted_text = extract_text_from_excel(content_bytes)
            elif file_metadata['mimeType'] == 'application/vnd.google-apps.presentation':
                request = drive_service.files().export_media(fileId=file_id, mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation')
                file_content = io.BytesIO()
                downloader = MediaIoBaseDownload(file_content, request)

                done = False
                while not done:
                    status, done = downloader.next_chunk()

                content_bytes = file_content.getvalue()
                extracted_text = extract_text_from_presentation(content_bytes)
            else:
                return {
                    "error": True,
                    "message": f"Unsupported Google Docs file type: {file_metadata['mimeType']}"
                }
        else:
            request = drive_service.files().get_media(fileId=file_id)
            file_content = io.BytesIO()
            downloader = MediaIoBaseDownload(file_content, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
                
            content_bytes = file_content.getvalue()
            extracted_text = extract_text_from_file(content_bytes, file_metadata['mimeType'], file_metadata['name'])

        document_cache[file_id] = (extracted_text, file_metadata)

        result = {
            "name": file_metadata['name'],
            "mime_type": file_metadata['mimeType'],
            "content": extracted_text[:max_chars],
            "truncated": len(extracted_text) > max_chars,
            "total_length": len(extracted_text),
            "from_cache": False
        }
        
        if include_metadata:
            result["metadata"] = file_metadata
            
        return result
    except Exception as e:
        logger.error(f"Error getting file content: {str(e)}")
        return {"error": True, "message": f"Error getting file content: {str(e)}"}

@mcp.tool()
async def search_drive_files(query: str, folder_id: str = None):
    """
    Searches for files in Google Drive that match a full-text search query.
    
    This tool performs a full-text search across file content and metadata in Google Drive.
    The search can be restricted to a specific folder if needed.
    
    Parameters:
        query (str): The text to search for within files.
                    Can be any string that might appear in file content or metadata.
                    For best results, use specific terms related to the content you're looking for.
        folder_id (str, optional): ID of a folder to restrict the search within.
                                  If provided, only searches within this folder (not recursive).
                                  If not provided, searches across all accessible files.
    
    Returns:
        list: A list of file objects that match the search query, each containing:
              - id: The unique Google Drive ID of the file
              - name: The display name of the file
              - mimeType: The MIME type of the file
              - createdTime: Timestamp when the file was created
              - modifiedTime: Timestamp when the file was last modified
        dict: Error response if the operation fails, with structure:
              {"error": True, "message": "Error message"}
    """
    try:
        drive_service = get_drive_service()

        search_query = f"fullText contains '{query}' and trashed=false"
        if folder_id:
            search_query += f" and '{folder_id}' in parents"

        response = drive_service.files().list(
            q=search_query,
            spaces='drive',
            fields='files(id,name,mimeType,createdTime,modifiedTime)'
        ).execute()

        files = response.get('files', [])
        return files
    except Exception as e:
        logger.error(f"Error searching files: {str(e)}")
        return {"error": True, "message": f"Error searching files: {str(e)}"}

async def verify_base_folder():
    try:
        drive_service = get_drive_service()
        folder = drive_service.files().get(
            fileId=SERVER_CONFIG["base_folder_id"],
            fields='id,name,mimeType'
        ).execute()

        if folder['mimeType'] != 'application/vnd.google-apps.folder':
            raise ValueError("Specified base_folder_id is not a folder")

        logger.info(f"Successfully verified base folder: {folder['name']} ({folder['id']})")
        SERVER_CONFIG["base_folder_name"] = folder['name']
        return True
    except Exception as e:
        logger.error(f"Error verifying base folder: {str(e)}")
        return False

async def initialize_server():
    logger.info("Initializing Google Drive Competitor Analysis MCP server...")

    service_account_file = os.environ.get('GOOGLE_SERVICE_ACCOUNT_FILE')
    if not service_account_file:
        logger.warning("Google Drive credentials not found in environment variables")
        return False

    if not os.path.exists(service_account_file):
        logger.error(f"Credentials file not found at: {service_account_file}")
        return False

    try:
        get_drive_service()
        logger.info("Successfully validated Google Drive credentials")

        if not SERVER_CONFIG.get("base_folder_id"):
            logger.error("Base folder ID not configured")
            return False

        if not await verify_base_folder():
            return False

    except Exception as e:
        logger.error(f"Failed to validate Google Drive credentials: {str(e)}")
        return False

    logger.info("Server initialized successfully")
    return True

async def shutdown_server():
    logger.info("Shutting down Google Drive Competitor Analysis MCP server...")
    document_cache.clear()
    logger.info("Server shutdown complete")

    
@mcp.tool()
async def list_all_files_recursive(folder_id: str = "", file_types: str = "all", max_depth: int = 10):
    """
    Recursively lists all files within a specified Google Drive folder and all its subfolders.
    
    This tool traverses the entire folder hierarchy starting from the specified parent folder,
    retrieving all files that match the specified file types, no matter how deeply nested.
    
    Parameters:
        folder_id (str, optional): The ID of the parent folder to start listing from.
                                  If not provided or empty, the base folder will be used.
        file_types (str, optional): Comma-separated list of file types to filter by.
                                   Default is "all" which returns all file types.
        max_depth (int, optional): Maximum depth of folder traversal. Default is 10 levels.
    
    Returns:
        list: A list of file objects with folder path information, each containing:
              - id: The unique Google Drive ID of the file
              - name: The display name of the file
              - mimeType: The MIME type of the file
              - folderPath: The path of folders from the root to this file
              - createdTime: Timestamp when the file was created
              - modifiedTime: Timestamp when the file was last modified
    """
    try:
        drive_service = get_drive_service()
        actual_folder = folder_id if folder_id else SERVER_CONFIG["base_folder_id"]
        
        all_files = []
        
        visited_folders = set()
        
        async def traverse_folder(current_folder_id, current_path, current_depth=0):
            if current_folder_id in visited_folders:
                return
            
            visited_folders.add(current_folder_id)
            
            if max_depth is not None and current_depth > max_depth:
                return
            
            try:
                folder_details = drive_service.files().get(fileId=current_folder_id, fields='name').execute()
                folder_name = folder_details.get('name', 'Unknown Folder')
                
                if current_path:
                    new_path = f"{current_path}/{folder_name}"
                else:
                    new_path = folder_name
            except Exception:
                new_path = current_path
            
            query = f"'{current_folder_id}' in parents and trashed=false"
            
            if file_types.lower() != "all":
                type_queries = []
                types = [t.strip().lower() for t in file_types.split(",")]

                for file_type in types:
                    if file_type == "pdf":
                        type_queries.append("mimeType='application/pdf'")
                    elif file_type in ["excel", "xlsx", "xls"]:
                        type_queries.append("mimeType='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'")
                        type_queries.append("mimeType='application/vnd.ms-excel'")
                        type_queries.append("mimeType='application/vnd.google-apps.spreadsheet'")
                    elif file_type in ["word", "doc", "docx"]:
                        type_queries.append("mimeType='application/vnd.openxmlformats-officedocument.wordprocessingml.document'")
                        type_queries.append("mimeType='application/msword'")
                        type_queries.append("mimeType='application/vnd.google-apps.document'")
                    elif file_type in ["ppt", "pptx", "presentation"]:
                        type_queries.append("mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'")
                        type_queries.append("mimeType='application/vnd.ms-powerpoint'")
                        type_queries.append("mimeType='application/vnd.google-apps.presentation'")
                    elif file_type in ["txt", "text"]:
                        type_queries.append("mimeType='text/plain'")

                if type_queries:
                    query += " and (" + " or ".join(type_queries) + ")"
            
            response = drive_service.files().list(
                q=query,
                spaces='drive',
                fields='files(id, name, mimeType, createdTime, modifiedTime, size)'
            ).execute()
            
            files = response.get('files', [])
            for file in files:
                if file['mimeType'] == 'application/vnd.google-apps.folder':
                    await traverse_folder(file['id'], new_path, current_depth + 1)
                else:
                    file['folderPath'] = new_path
                    all_files.append(file)
        
        await traverse_folder(actual_folder, "")
        
        return all_files
    except Exception as e:
        logger.error(f"Error recursively listing files: {str(e)}")
        return {"error": True, "message": f"Error recursively listing files: {str(e)}"}
@mcp.tool()
async def search_all_folders(query: str, folder_id: str = "", max_depth: int = 10):
    """
    Searches for files across all folders and subfolders that match a full-text search query.
    
    This tool performs a comprehensive search across the entire folder hierarchy,
    starting from the specified folder, looking for content matching the query.
    
    Parameters:
        query (str): The text to search for within files.
                    Can be any string that might appear in file content or metadata.
        folder_id (str, optional): ID of a folder to restrict the search within.
                                  If not provided or empty, searches from the base folder.
        max_depth (int, optional): Maximum depth of folder traversal. Default is 10 levels.
    
    Returns:
        list: A list of file objects that match the search query, each containing:
              - id: The unique Google Drive ID of the file
              - name: The display name of the file
              - mimeType: The MIME type of the file
              - folderPath: The path to the file
              - snippet: A text snippet showing the context of the match
              - createdTime: Timestamp when the file was created
              - modifiedTime: Timestamp when the file was last modified
    """
    try:
        drive_service = get_drive_service()
        root_folder = folder_id if folder_id else SERVER_CONFIG["base_folder_id"]
        
        all_files = await list_all_files_recursive(root_folder, "all", max_depth)
        
        if isinstance(all_files, dict) and all_files.get("error"):
            return all_files
        
        search_results = []
        
        search_query = f"fullText contains '{query}' and trashed=false"
        
        response = drive_service.files().list(
            q=search_query,
            spaces='drive',
            fields='files(id,name,mimeType,createdTime,modifiedTime)'
        ).execute()
        
        found_file_ids = {file['id'] for file in response.get('files', [])}
        
        for file in all_files:
            if file['id'] in found_file_ids:
                try:
                    file_content = await get_file_content(file['id'], max_chars=5000)
                    if isinstance(file_content, dict) and not file_content.get("error"):
                        content = file_content.get("content", "")
                        query_pos = content.lower().find(query.lower())
                        if query_pos > -1:
                            start = max(0, query_pos - 100)
                            end = min(len(content), query_pos + len(query) + 100)
                            file['snippet'] = f"...{content[start:end]}..."
                        else:
                            file['snippet'] = "Match found, but no context available"
                    else:
                        file['snippet'] = "Content not accessible"
                except Exception:
                    file['snippet'] = "Unable to generate snippet"
                
                search_results.append(file)
        
        return search_results
    except Exception as e:
        logger.error(f"Error searching across folders: {str(e)}")
        return {"error": True, "message": f"Error searching across folders: {str(e)}"}
    
@mcp.tool()
async def analyze_documents(query: str, folder_id: str = "", max_files: int = 5, file_types: str = "all"):
    """
    Analyzes the content of multiple documents to find relevant information related to a query.
    
    This tool searches for documents matching the query, extracts their content,
    and provides a summary of the relevant information found across all analyzed documents.
    
    Parameters:
        query (str): The search query or question to analyze documents for.
        folder_id (str, optional): The ID of the folder to search in and its subfolders.
                                   If not provided or empty, the base folder will be used.
        max_files (int, optional): Maximum number of files to analyze. Default is 5.
        file_types (str, optional): Comma-separated list of file types to consider.
                                   Default is "all" which includes all file types.
    
    Returns:
        dict: A dictionary containing:
              - query: The original query
              - file_count: Number of files analyzed
              - files: List of analyzed files with their metadata
              - summary: Summary of findings across all documents
    """
    try:
        actual_folder = folder_id if folder_id else SERVER_CONFIG["base_folder_id"]
        
        search_results = await search_drive_files(query, actual_folder)
        
        if isinstance(search_results, dict) and search_results.get("error"):
            return search_results
        
        files_to_analyze = search_results[:max_files]
        
        analyzed_files = []
        all_content = []
        
        for file in files_to_analyze:
            file_content = await get_file_content(file['id'])
            
            if isinstance(file_content, dict) and not file_content.get("error"):
                file_info = {
                    "id": file['id'],
                    "name": file['name'],
                    "path": file.get('folderPath', ''),
                    "mimeType": file['mimeType'],
                    "modifiedTime": file.get('modifiedTime', '')
                }
                
                content = file_content.get("content", "")
                
                if len(content) > 500:
                    file_info["content_preview"] = content[:500] + "..."
                else:
                    file_info["content_preview"] = content
                
                analyzed_files.append(file_info)
                
                file_header = f"Document: {file['name']}\n"
                all_content.append(file_header + content)
        
        combined_content = "\n\n---\n\n".join(all_content)
        
        query_terms = query.lower().split()
        summary = f"Analysis results for query: '{query}'\n\n"
        
        if not analyzed_files:
            summary += "No relevant documents found."
        else:
            summary += f"Analyzed {len(analyzed_files)} documents.\n\n"
            summary += "Key findings:\n"
            
            for i, file in enumerate(analyzed_files):
                file_content = await get_file_content(file['id'])
                if isinstance(file_content, dict) and not file_content.get("error"):
                    content = file_content.get("content", "")
                    
                    paragraphs = content.split('\n\n')
                    relevant_paragraphs = []
                    
                    for para in paragraphs:
                        if any(term in para.lower() for term in query_terms):
                            relevant_paragraphs.append(para)
                    
                    if relevant_paragraphs:
                        summary += f"\n- From {file['name']}:\n"
                        
                        for j, para in enumerate(relevant_paragraphs[:3]):
                            if len(para) > 300:
                                para = para[:300] + "..."
                            summary += f"  {j+1}. {para}\n"
                        
                        if len(relevant_paragraphs) > 3:
                            summary += f"  ... and {len(relevant_paragraphs) - 3} more relevant sections\n"
        
        return {
            "query": query,
            "file_count": len(analyzed_files),
            "files": analyzed_files,
            "summary": summary
        }
    except Exception as e:
        logger.error(f"Error analyzing documents: {str(e)}")
        return {"error": True, "message": f"Error analyzing documents: {str(e)}"}

@mcp.tool()
async def simple_document_analysis(query: str, folder_id: str = None, max_files: int = 5):
    """
    Performs a simple analysis of documents matching a query, with minimal dependencies.
    
    This tool directly searches Google Drive for files matching the query
    and extracts their content for analysis.
    
    Parameters:
        query (str): The search query or question to analyze documents for.
        folder_id (str, optional): The ID of the folder to search in.
                                  If not provided, the base folder will be used.
        max_files (int, optional): Maximum number of files to analyze. Default is 5.
    
    Returns:
        dict: Analysis results and document information
    """
    try:
        drive_service = get_drive_service()
        actual_folder = folder_id if folder_id else SERVER_CONFIG["base_folder_id"]
        
        search_query = f"fullText contains '{query}' and trashed=false"
        if folder_id:
            search_query += f" and '{actual_folder}' in parents"
        
        response = drive_service.files().list(
            q=search_query,
            spaces='drive',
            fields='files(id,name,mimeType,createdTime,modifiedTime)',
            pageSize=max_files
        ).execute()
        
        files = response.get('files', [])
        
        if not files:
            return {
                "query": query,
                "message": "No files found matching your query.",
                "files": []
            }
        
        results = []
        
        for file in files:
            content_result = await get_file_content(file['id'])
            
            file_data = {
                "id": file['id'],
                "name": file['name'],
                "type": file['mimeType'],
                "modified": file.get('modifiedTime', 'Unknown')
            }
            
            if isinstance(content_result, dict) and not content_result.get("error"):
                content = content_result.get("content", "")
                
                snippets = []
                query_terms = query.lower().split()
                
                if isinstance(content, str):
                    paragraphs = content.split('\n\n')
                    
                    for para in paragraphs:
                        if para.strip() and any(term in para.lower() for term in query_terms):
                            if len(para) > 300:
                                para = para[:300] + "..."
                            snippets.append(para.strip())
                    
                    if len(snippets) > 3:
                        snippets = snippets[:3]
                
                file_data["relevant_content"] = snippets
                file_data["content_available"] = True
            else:
                error_msg = content_result.get("message", "Unknown error") if isinstance(content_result, dict) else "Failed to retrieve content"
                file_data["error"] = error_msg
                file_data["content_available"] = False
            
            results.append(file_data)
        
        summary = f"Analysis for query: '{query}'\n\n"
        summary += f"Found {len(results)} relevant documents:\n\n"
        
        for i, file in enumerate(results):
            summary += f"{i+1}. {file['name']}\n"
            
            if file.get("content_available"):
                if file.get("relevant_content"):
                    summary += "   Relevant excerpts:\n"
                    for j, snippet in enumerate(file.get("relevant_content", [])):
                        summary += f"   - {snippet}\n\n"
                else:
                    summary += "   No highly relevant content found in this document.\n\n"
            else:
                summary += f"   Could not analyze content: {file.get('error', 'Unknown error')}\n\n"
        
        return {
            "query": query,
            "file_count": len(results),
            "files": results,
            "summary": summary
        }
        
    except Exception as e:
        logger.error(f"Error in simple_document_analysis: {str(e)}")
        import traceback
        logger.error(traceback.format_exc())
        return {
            "error": True,
            "message": f"Error analyzing documents: {str(e)}",
            "query": query
        }


@mcp.tool()
async def clear_document_cache():
    """
    Clears the document cache to free up memory and force fresh document retrieval.
    
    This can be useful if documents have been updated in Google Drive and you want
    to ensure you're seeing the latest content.
    
    Returns:
        dict: Result of the cache clearing operation
    """
    try:
        cache_size = len(document_cache)
        document_cache.clear()
        return {
            "success": True,
            "message": f"Successfully cleared document cache ({cache_size} items removed)"
        }
    except Exception as e:
        logger.error(f"Error clearing document cache: {str(e)}")
        return {
            "error": True,
            "message": f"Error clearing document cache: {str(e)}"
        }
@mcp.tool()
async def get_folder_structure(folder_id: str = "", max_depth: int = 5):
    """
    Gets the hierarchical folder structure starting from a specified folder.
    
    This tool returns the complete folder hierarchy, showing how folders are nested
    within each other, which is useful for understanding the organization of opportunity folders.
    
    Parameters:
        folder_id (str, optional): The ID of the parent folder to start mapping from.
                                  If not provided or empty, the base folder will be used.
        max_depth (int, optional): Maximum depth of folder traversal. Default is 5 levels.
    
    Returns:
        dict: A nested dictionary representing the folder structure, where:
              - Each key is a folder name
              - Each value is either:
                - A dictionary of subfolders (recursive structure)
                - A list of file metadata if the folder contains files
    """
    try:
        drive_service = get_drive_service()
        actual_folder = folder_id if folder_id else SERVER_CONFIG["base_folder_id"]
        
        visited_folders = set()
        
        async def build_structure(current_folder_id, current_depth=0):
            if current_folder_id in visited_folders:
                return {"error": "Cyclic folder reference detected"}
            
            visited_folders.add(current_folder_id)
            
            if max_depth is not None and current_depth > max_depth:
                return {"truncated": "Max depth reached"}
            
            try:
                folder_details = drive_service.files().get(
                    fileId=current_folder_id, 
                    fields='name'
                ).execute()
                folder_name = folder_details.get('name', 'Unknown Folder')
            except Exception as e:
                return {"error": f"Unable to access folder: {str(e)}"}
            
            structure = {
                "name": folder_name,
                "id": current_folder_id
            }
            
            query = f"'{current_folder_id}' in parents and trashed=false"
            response = drive_service.files().list(
                q=query,
                spaces='drive',
                fields='files(id, name, mimeType)'
            ).execute()
            
            items = response.get('files', [])
            
            subfolders = []
            files = []
            
            for item in items:
                if item['mimeType'] == 'application/vnd.google-apps.folder':
                    subfolder_structure = await build_structure(item['id'], current_depth + 1)
                    subfolders.append(subfolder_structure)
                else:
                    files.append({
                        "id": item['id'],
                        "name": item['name'],
                        "mimeType": item['mimeType']
                    })
            
            if subfolders:
                structure["subfolders"] = subfolders
            
            if files:
                structure["files"] = files
            
            return structure
        
        folder_structure = await build_structure(actual_folder)
        
        return folder_structure
    except Exception as e:
        logger.error(f"Error getting folder structure: {str(e)}")
        return {"error": True, "message": f"Error getting folder structure: {str(e)}"}

if __name__ == "__main__":
    print("Starting Google Drive Competitor Analysis MCP server...")
    loop = asyncio.get_event_loop()
    try:
        loop.run_until_complete(initialize_server())
        mcp.run(transport='stdio')
        
    except Exception as e:
        logger.error(f"Error starting server: {e}")
        print(f"Error starting server: {e}")

    finally:
        try:
            loop.run_until_complete(shutdown_server())
        except Exception as e:
            logger.error(f"Error during server shutdown: {e}")